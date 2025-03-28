import os
import asyncio
import pickle
import concurrent.futures
import streamlit as st
import docx
import pptx
import pytesseract
import pandas as pd
from PIL import Image
from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
from langchain_community.document_loaders import PyPDFLoader
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import SystemMessage, HumanMessage

# ✅ Load environment variables
load_dotenv()

pinecone_api_key = os.getenv("PINECONE_API_KEY")
openai_api_key = os.getenv("OPENAI_API_KEY")

# ✅ Initialize Pinecone
pc = Pinecone(api_key=pinecone_api_key)

database_index_name = "database-index"
brd_index_name = "brd-index"

# Ensure indexes exist
for index_name in [database_index_name, brd_index_name]:
    if index_name not in pc.list_indexes().names():
        pc.create_index(
            name=index_name,
            dimension=1536,  # OpenAI embedding dimension
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1")
        )

# ✅ Initialize Pinecone index objects
db_index = pc.Index(database_index_name)
brd_index = pc.Index(brd_index_name)

# ✅ Initialize OpenAI Model and Embeddings
embedding_model = OpenAIEmbeddings(api_key=openai_api_key)
llm = ChatOpenAI(model_name="gpt-3.5-turbo", temperature=0.3, api_key=openai_api_key)

# ✅ Streamlit UI
st.set_page_config(page_title="AI Document Query", layout="wide")
st.title("📄 AI-Powered Document Query System")

# ✅ Upload Files
uploaded_db_file = st.file_uploader("📂 Upload Database File", type=["pdf", "docx", "pptx", "txt", "xls", "xlsx", "png", "jpg", "jpeg"])
uploaded_brd_file = st.file_uploader("📂 Upload BRD Document", type=["pdf", "docx", "pptx", "txt", "xls", "xlsx", "png", "jpg", "jpeg"])

# ✅ Cache Embeddings
def cache_embeddings(file_path, embeddings):
    cache_path = f"{file_path}.pkl"
    with open(cache_path, "wb") as f:
        pickle.dump(embeddings, f)

def load_cached_embeddings(file_path):
    cache_path = f"{file_path}.pkl"
    if os.path.exists(cache_path):
        with open(cache_path, "rb") as f:
            return pickle.load(f)
    return None

# ✅ Extract Text Function
async def extract_text_async(file_path):
    loop = asyncio.get_event_loop()
    with concurrent.futures.ThreadPoolExecutor() as pool:
        return await loop.run_in_executor(pool, lambda: extract_text_from_file(file_path))

def extract_text_from_file(file_path):
    ext = file_path.split('.')[-1].lower()
    if ext == "pdf":
        loader = PyPDFLoader(file_path)
        return "\n".join([doc.page_content.strip() for doc in loader.load()])
    elif ext in ["doc", "docx"]:
        doc = docx.Document(file_path)
        return "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
    elif ext in ["ppt", "pptx"]:
        presentation = pptx.Presentation(file_path)
        return "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext in ["xls", "xlsx"]:
        df_dict = pd.read_excel(file_path, sheet_name=None)
        return "\n".join([df.to_markdown(index=False) for df in df_dict.values()])
    elif ext in ["png", "jpg", "jpeg"]:
        return pytesseract.image_to_string(Image.open(file_path)).strip()
    return "⚠️ Unsupported file type."

# ✅ Batched Embedding and Upsert
def batch_embed_and_upsert(index, file_path, text_chunks, prefix, batch_size=32):
    embeddings = embedding_model.embed_documents(text_chunks)
    cache_embeddings(file_path, embeddings)
    for i in range(0, len(text_chunks), batch_size):
        records = [
            (f"{prefix}-{i + j}", embeddings[i + j], {"text": text_chunks[i + j]})
            for j in range(min(batch_size, len(text_chunks) - i))
        ]
        index.upsert(records)

# ✅ Process Uploads
async def process_upload(uploaded_file, index, prefix):
    if uploaded_file is not None:
        file_path = f"./uploaded_files/{uploaded_file.name}"
        os.makedirs("uploaded_files", exist_ok=True)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        cached = load_cached_embeddings(file_path)
        if cached is None:
            text_data = await extract_text_async(file_path)
            splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            text_chunks = splitter.split_text(text_data)
            batch_embed_and_upsert(index, file_path, text_chunks, prefix)
            st.success(f"✅ Stored embeddings for '{uploaded_file.name}'")
        else:
            st.info(f"⚡ Using cached embeddings for '{uploaded_file.name}'")

# ✅ Main Execution Loop
async def main():
    await asyncio.gather(
        process_upload(uploaded_db_file, db_index, "database"),
        process_upload(uploaded_brd_file, brd_index, "brd")
    )

    query = st.text_input("💬 Ask a Question")

    if st.button("🔍 Get Answer") and query:
        query_embedding = embedding_model.embed_query(query)

        db_results = db_index.query(vector=query_embedding, top_k=5, include_metadata=True)
        brd_results = brd_index.query(vector=query_embedding, top_k=5, include_metadata=True)

        retrieved_texts = "\n".join([match['metadata']['text'] for match in db_results['matches'] + brd_results['matches']])

        final_prompt = [
    SystemMessage(content=(
        "You are an advanced AI assistant specializing in analyzing Business Requirement Documents (BRD) and comparing them with database information. "
        "Your task is to thoroughly analyze the BRD, extract component-related requirements, and identify matching or reusable components from the database. "
        "Your goal is to deliver a structured and detailed comparison, highlighting component details and their compatibility with the BRD requirements."
    )),
    HumanMessage(content=(
        f"### Instructions:\n"
        f"1. **Parse and Extract Requirements**:\n"
        f"   - Thoroughly analyze the provided BRD document to identify and extract all component or code-related requirements.\n\n"
        f"2. **Compare with Database**:\n"
        f"   - Search through the provided database to find components that strictly match or can be reused for the identified BRD requirements.\n"
        f"   - Ensure a detailed comparison by validating component compatibility, usage, and any discrepancies between the BRD and the database.\n\n"
        f"3. **Provide Detailed Information**:\n"
        f"   - For **each matched or related component**, deliver the following details in a structured format:\n"
        f"      - **Component Name**: The exact name of the component.\n"
        f"      - **Description**: A brief but clear explanation of the component's functionality.\n"
        f"      - **Projects Used In**: A list of projects where this component is currently used (include all associated projects).\n"
        f"      - **Compatibility Details**: Explain whether and how the component can be reused for the BRD requirement. Highlight compatibility issues if any.\n"
        f"   - If a component is **not found** or is **incompatible**, explicitly mention it and provide potential alternatives if available.\n\n"
        f"4. **Handle Ambiguities**:\n"
        f"   - If the information is incomplete, ambiguous, or requires further clarification, clearly indicate this while providing the best possible analysis.\n\n"
        f"### Context Information:\n{retrieved_texts}\n\n"
        f"### User Query:\n{query}\n\n"
        f"### Response Format:\n"
        f"Deliver the output in a clean, well-structured JSON-like format as shown below:\n\n"
        f"[\n"
        f"  {{\n"
        f"    'Component Name': 'Component A',\n"
        f"    'Description': 'Handles user authentication and session management.',\n"
        f"    'Projects Used In': ['Project X', 'Project Y'],\n"
        f"    'Compatibility Details': 'Fully compatible with BRD requirements for user authentication. Requires minor modification for multi-factor authentication.'\n"
        f"  }},\n"
        f"  {{\n"
        f"    'Component Name': 'Component B',\n"
        f"    'Description': 'Processes financial transactions.',\n"
        f"    'Projects Used In': ['Project Z'],\n"
        f"    'Compatibility Details': 'Not compatible with the BRD due to outdated payment protocols. Suggest using Component C as an alternative.'\n"
        f"  }}\n"
        f"]\n\n"
        f"Ensure the response is **detailed**, **accurate**, and **directly addresses** the BRD requirements."
    ))
]


        response = llm.invoke(final_prompt).content

        st.subheader("💡 AI Response")
        st.text_area("AI Response", response, height=300)

asyncio.run(main())